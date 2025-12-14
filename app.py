import streamlit as st
import streamlit.components.v1 as components
from openai import OpenAI
from dotenv import load_dotenv
from pypdf import PdfReader
from pptx import Presentation
from docx import Document
import io
import os
import json   # ← Add this line

# -------------------- LOAD API --------------------
load_dotenv()
client = OpenAI(api_key=os.getenv("OPENAI_API_KEY"))

# -------------------- PAGE CONFIG --------------------
st.set_page_config(
    page_title="Python Data Structures Tutorial",
    layout="wide",
    initial_sidebar_state="expanded",
    page_icon="📚"
)

# -------------------- CSS --------------------
st.markdown("""
<style>
@import url('https://fonts.googleapis.com/css2?family=Inter:wght@300;400;500;600;700&display=swap');
body {background-color:#ffffff;color:#1a1a1a;font-family:'Inter',sans-serif;}
.stApp {background-color:#ffffff;}
.block-container {padding-top:2rem;max-width:1200px;margin:0 auto;}
.doc-title {font-size:48px;font-weight:700;color:#2d3748;text-align:center;}
.doc-subtitle {font-size:20px;color:#718096;text-align:center;margin-bottom:3rem;}
.doc-box {background:#fff;padding:30px;border-radius:12px;
box-shadow:0 4px 6px -1px rgba(0,0,0,.1);}
.stSidebar {background-color:#f7fafc;}
.copy-btn button {
    width:100%;
    padding:12px;
    font-size:16px;
    background:#2563eb;
    color:white;
    border:none;
    border-radius:8px;
    cursor:pointer;
}
.copy-btn button:hover {background:#1d4ed8;}
.example-code {background-color: #f7fafc; padding: 15px; border-radius: 8px; border-left: 4px solid #2563eb;}
</style>
""", unsafe_allow_html=True)

# -------------------- HEADER --------------------
st.markdown('<div class="doc-title">Python Data Structures Tutorial</div>', unsafe_allow_html=True)
st.markdown('<div class="doc-subtitle">Comprehensive Guide with Examples and Best Practices</div>', unsafe_allow_html=True)

# -------------------- SESSION STATE --------------------
if "doc_context" not in st.session_state:
    st.session_state.doc_context = ""

if "generated_code" not in st.session_state:
    st.session_state.generated_code = ""

# -------------------- SIDEBAR --------------------
with st.sidebar:
    st.title("📘 Topics")
    topic = st.radio(
        "Navigate Topics",
        [
            "Introduction", "Lists", "Tuples", "Dictionaries",
            "Sets", "Strings", "Advanced: Stacks & Queues",
            "Advanced: Trees & Graphs"
        ]
    )

    st.markdown("---")
    
    st.markdown("<br><br><br><br><br><br><br><br><br><br><br><br><br><br><br><br><br><br><br>", unsafe_allow_html=True)

    with st.expander("Help Desk", expanded=False):
        uploaded = st.file_uploader(
            "",
            type=["pdf", "pptx", "docx"]
        )

        if uploaded:
            if uploaded.size > 500 * 1024 * 1024:
                st.error("File too large (max 500MB)")
            else:
                extracted = ""
                data = uploaded.read()

                if uploaded.name.endswith(".pdf"):
                    reader = PdfReader(io.BytesIO(data))
                    for p in reader.pages:
                        text = p.extract_text()
                        if text:
                            extracted += text + "\n"

                elif uploaded.name.endswith(".pptx"):
                    pres = Presentation(io.BytesIO(data))
                    for slide in pres.slides:
                        for shape in slide.shapes:
                            if hasattr(shape, "text"):
                                extracted += shape.text + "\n"

                elif uploaded.name.endswith(".docx"):
                    doc = Document(io.BytesIO(data))
                    for para in doc.paragraphs:
                        extracted += para.text + "\n"

                st.session_state.doc_context = extracted

        user_prompt = st.text_input(
            "Your Query",
            placeholder="Send query here..."
        )

        if st.button("Submit"):
            if not user_prompt:
                st.warning("Please enter a querry")
            else:
                final_prompt = f"""
IMPORTANT:
This document is a STRICT SPECIFICATION, not guidance.
Any deviation will FAIL validation.

GENERAL RULES:
- Output ONLY valid, executable Python code
- EXACTLY ONE short header comment on the first line
- No explanations, markdown, or backticks
- Use ONLY Python standard library
- Code must run as-is in Python 3

WHEN CORRECTING EXISTING CODE:
- Preserve ALL variable, function, and class names
- Do NOT rename identifiers
- Modify logic ONLY where required

WHEN ASKED FOR PARTIAL IMPLEMENTATION:
- Implement ONLY requested milestones
- Do NOT remove previous functionality
- Do NOT implement future milestones

ENGINEERING RULES:
- Simple, explicit logic
- No clever tricks
- Handle edge cases
- No partial implementations
- No TODO or pass

DOCUMENT CONTEXT:
{st.session_state.doc_context[:6000]}

PREVIOUS CODE:
{st.session_state.generated_code}

USER REQUEST:
{user_prompt}
"""

                with st.spinner("Sending querry..."):
                    response = client.chat.completions.create(
                        model="gpt-4o",
                        messages=[
                            {"role": "system", "content": "Return ONLY valid executable Python code."},
                            {"role": "user", "content": final_prompt}
                        ],
                        temperature=0.0,
                        max_tokens=5000
                    )

                raw_output = response.choices[0].message.content

                lines = []
                for line in raw_output.splitlines():
                    s = line.strip()
                    if not s:
                        continue
                    if s.startswith(("```", "---")):
                        continue
                    lines.append(line)

                code = "\n".join(lines).strip()

                try:
                    compile(code, "<generated>", "exec")
                    if "TODO" in code or "\npass\n" in code:
                        raise SyntaxError
                    st.session_state.generated_code = code
                    st.success("successfull")
                except SyntaxError:
                    st.error("Failed validation")

                # -------- COPY TO CLIPBOARD (WORKING 100%) --------
        if st.session_state.generated_code:
            st.markdown("### ✅ Generated Code")
            st.code(st.session_state.generated_code, language="python")

            # This is the working copy button
            components.html(
                f"""
                <script>
                function copyCode() {{
                    navigator.clipboard.writeText({json.dumps(st.session_state.generated_code)});
                    const btn = document.getElementById('copyBtn');
                    btn.innerHTML = '✅ Copied!';
                    setTimeout(() => btn.innerHTML = '📋 Copy Code', 2000);
                }}
                </script>
                <div class="copy-btn">
                    <button id="copyBtn" onclick="copyCode()">
                        📋 Copy Code
                    </button>
                </div>
                """,
                height=100
            )

# -------------------- MAIN CONTENT --------------------
st.markdown('<div class="doc-box">', unsafe_allow_html=True)

if topic == "Introduction":
    st.markdown("""
    # Introduction to Python Data Structures
    
    Data structures are fundamental building blocks in programming that allow us to organize, store, and manage data efficiently. In Python, data structures are essential for solving complex problems, from simple list management to advanced algorithmic implementations.
    
    ## Why Data Structures Matter
    - **Efficiency**: Choosing the right data structure can significantly improve time and space complexity (e.g., O(1) lookups with dictionaries vs. O(n) with lists).
    - **Readability**: Well-structured data makes code cleaner and easier to maintain.
    - **Scalability**: As applications grow, proper data structures prevent performance bottlenecks.
    
    Python provides several built-in data structures that are versatile and powerful:
    - **Lists**: Ordered, mutable collections.
    - **Tuples**: Ordered, immutable collections.
    - **Dictionaries**: Unordered key-value pairs.
    - **Sets**: Unordered, unique elements.
    - **Strings**: Immutable sequences of characters.
    
    We'll also cover advanced structures like stacks, queues, trees, and graphs, which can be implemented using Python's primitives.
    
    ## Basic Principles
    - **Mutability**: Can the structure be changed after creation? (Lists: yes; Tuples: no)
    - **Order**: Are elements in a specific sequence? (Lists: yes; Sets: no)
    - **Uniqueness**: Must elements be unique? (Sets: yes; Lists: no)
    
    Let's dive deeper into each with hands-on examples!
    """)
    
    col1, col2 = st.columns(2)
    with col1:
        st.markdown("### Quick Tip: Time Complexities")
        st.markdown("""
        | Operation | List | Dict | Set |
        |-----------|------|------|-----|
        | Access    | O(1) | O(1) | O(1)|
        | Insert    | O(1)*| O(1)| O(1)|
        | Delete    | O(n) | O(1)| O(1)|
        *Amortized
        """)
    with col2:
        st.markdown("""
        ### Example: Simple Data Organization
        Imagine tracking student grades:
        - List: `[85, 92, 78]`
        - Dict: `{'Alice': 85, 'Bob': 92}`
        """)
        st.code("""
        grades = [85, 92, 78]
        print(sum(grades) / len(grades))  # Average: 85.0
        """, language="python")

elif topic == "Lists":
    st.markdown("""
    # Lists in Python
    
    Lists are one of the most versatile data structures in Python. They are ordered, mutable sequences that can hold elements of any type. Lists are defined using square brackets `[]` and are zero-indexed.
    
    ## Creating Lists
    - Empty list: `my_list = []`
    - With elements: `fruits = ['apple', 'banana', 'cherry']`
    - From iterable: `numbers = list(range(5))  # [0, 1, 2, 3, 4]`
    
    ## Accessing Elements
    - Indexing: `fruits[0]` → `'apple'`
    - Negative indexing: `fruits[-1]` → `'cherry'`
    - Slicing: `fruits[1:3]` → `['banana', 'cherry']`
    
    ## Common Methods
    - `append(x)`: Add to end. O(1)
    - `extend(iterable)`: Add multiple. O(k)
    - `insert(i, x)`: Insert at index. O(n)
    - `remove(x)`: Remove first occurrence. O(n)
    - `pop(i)`: Remove and return by index. O(n) or O(1) for end
    - `clear()`: Remove all
    - `index(x)`: Find index. O(n)
    - `count(x)`: Count occurrences. O(n)
    - `sort()`: In-place sort. O(n log n)
    - `reverse()`: In-place reverse. O(n)
    
    ## List Comprehensions
    Powerful for creating lists: `[x**2 for x in range(10) if x % 2 == 0]`
    
    ## Examples
    """)
    
    st.subheader("Example 1: Basic Operations")
    st.markdown("Building a shopping list.")
    st.code("""
    shopping = []
    shopping.append('milk')
    shopping.extend(['bread', 'eggs'])
    shopping.insert(1, 'butter')
    print(shopping)  # ['milk', 'butter', 'bread', 'eggs']
    shopping.remove('bread')
    print(shopping[-1])  # 'eggs'
    """, language="python")
    
    st.subheader("Example 2: List Comprehension for Filtering")
    st.markdown("Generate even squares up to 20.")
    st.code("""
    evens_squared = [x**2 for x in range(1, 11) if x % 2 == 0]
    print(evens_squared)  # [4, 16, 36, 64, 100]
    """, language="python")
    
    st.subheader("Example 3: Sorting and Reversing")
    st.markdown("Sort a list of numbers and reverse it.")
    st.code("""
    numbers = [3, 1, 4, 1, 5, 9]
    numbers.sort()
    print(numbers)  # [1, 1, 3, 4, 5, 9]
    numbers.reverse()
    print(numbers)  # [9, 5, 4, 3, 1, 1]
    """, language="python")
    
    st.subheader("Example 4: Nested Lists (Matrices)")
    st.markdown("Represent a 2D grid.")
    st.code("""
    matrix = [[1, 2, 3], [4, 5, 6], [7, 8, 9]]
    print(matrix[1][2])  # 6 (row 1, col 2)
    # Flatten
    flattened = [item for row in matrix for item in row]
    print(flattened)  # [1, 2, 3, 4, 5, 6, 7, 8, 9]
    """, language="python")
    
    st.subheader("Best Practices")
    st.markdown("""
    - Use lists for ordered, changeable collections.
    - Avoid frequent inserts/deletes in middle (use deques for queues).
    - For large lists, consider slicing carefully to avoid O(n) copies.
    """)
    
    col1, col2 = st.columns(2)
    with col1:
        st.markdown("### Edge Cases")
        st.code("""
        empty = []
        print(len(empty))  # 0
        print(empty[0])  # IndexError!
        """, language="python")
    with col2:
        st.markdown("### Performance Tip")
        st.write("Use `append` over `+` for concatenation in loops.")

elif topic == "Tuples":
    st.markdown("""
    # Tuples in Python
    
    Tuples are immutable, ordered sequences similar to lists but cannot be modified after creation. They are defined with parentheses `()` and are useful for fixed data, function returns, and as dictionary keys.
    
    ## Creating Tuples
    - Empty: `t = ()`
    - Single element: `t = (42,)` (comma needed!)
    - Multiple: `point = (3, 4)`
    - From list: `t = tuple([1, 2, 3])`
    
    ## Accessing Elements
    Same as lists: indexing, slicing, negative indices.
    
    ## Why Use Tuples?
    - **Immutability**: Safer for constants, prevents accidental changes.
    - **Hashable**: Can be dictionary keys or set elements.
    - **Memory Efficient**: Slightly smaller than lists.
    - **Unpacking**: `x, y = (3, 4)` → x=3, y=4
    
    ## Common Operations
    - No methods like append/remove (immutable!).
    - `index(x)`, `count(x)` work.
    - Concatenation: `(1,2) + (3,4)` → `(1,2,3,4)`
    - Repetition: `(1,2) * 3` → `(1,2,1,2,1,2)`
    
    ## Examples
    """)
    
    st.subheader("Example 1: Coordinates")
    st.markdown("Storing 2D points.")
    st.code("""
    origin = (0, 0)
    print(origin[0] + origin[1])  # 0
    # origin[0] = 1  # TypeError: immutable!
    """, language="python")
    
    st.subheader("Example 2: Multiple Return Values")
    st.markdown("Functions returning tuples.")
    st.code("""
    def get_stats(numbers):
        return (sum(numbers), len(numbers), sum(numbers)/len(numbers))
    
    data = [1, 2, 3, 4, 5]
    total, count, avg = get_stats(data)
    print(f"Total: {total}, Avg: {avg}")  # Total: 15, Avg: 3.0
    """, language="python")
    
    st.subheader("Example 3: Tuple Unpacking")
    st.markdown("Swapping variables elegantly.")
    st.code("""
    a, b = 1, 2
    a, b = b, a  # Swap
    print(a, b)  # 2 1
    
    # Extended unpacking (Python 3.0+)
    numbers = (1, 2, 3, 4, 5)
    head, *tail = numbers
    print(head, tail)  # 1 [2, 3, 4, 5]
    """, language="python")
    
    st.subheader("Example 4: Nested Tuples")
    st.markdown("Representing a matrix of points.")
    st.code("""
    points = ((0,0), (1,2), (3,4))
    distances = [(x**2 + y**2)**0.5 for x,y in points]
    print(distances)  # [0.0, 2.236..., 5.0]
    """, language="python")
    
    st.subheader("Best Practices")
    st.markdown("""
    - Use tuples for heterogeneous data (e.g., records).
    - Prefer tuples over lists for fixed-size data.
    - Use namedtuples from `collections` for readable records.
    """)
    
    st.info("Tuples are faster for iteration than lists due to immutability.")

elif topic == "Dictionaries":
    st.markdown("""
    # Dictionaries in Python
    
    Dictionaries (dicts) are unordered (pre-Python 3.7), mutable collections of key-value pairs. Keys must be hashable (immutable types like strings, tuples). Defined with curly braces `{}`.
    
    ## Creating Dictionaries
    - Empty: `d = {}`
    - With elements: `person = {'name': 'Alice', 'age': 30}`
    - From iterables: `d = dict([('a',1), ('b',2)])`
    - `dict.fromkeys(keys, value)`: `{'a': None, 'b': None}`
    
    ## Accessing Elements
    - By key: `person['name']` → `'Alice'` (KeyError if missing)
    - `get(key, default)`: Safe access, `person.get('age', 0)` → 30
    - `keys()`, `values()`, `items()`: Views of contents
    
    ## Common Methods
    - `update(other)`: Merge dicts. O(len(other))
    - `pop(key)`: Remove and return value
    - `clear()`: Empty dict
    - `setdefault(key, default)`: Get or insert default
    - `copy()`: Shallow copy
    
    ## Dictionary Comprehensions
    `{k: v**2 for k,v in enumerate(range(5))}` → `{0:0, 1:1, 2:4, 3:9, 4:16}`
    
    ## Examples
    """)
    
    st.subheader("Example 1: Student Database")
    st.markdown("Storing grades by name.")
    st.code("""
    grades = {'Alice': 85, 'Bob': 92, 'Charlie': 78}
    print(grades['Bob'])  # 92
    grades['David'] = 88  # Add
    avg = sum(grades.values()) / len(grades)
    print(f"Average: {avg:.1f}")  # Average: 85.8
    """, language="python")
    
    st.subheader("Example 2: Safe Access with get()")
    st.markdown("Avoid KeyErrors.")
    st.code("""
    user = {'name': 'Eve', 'email': 'eve@example.com'}
    phone = user.get('phone', 'N/A')
    print(phone)  # N/A
    """, language="python")
    
    st.subheader("Example 3: Iterating Dictionaries")
    st.markdown("Loop over items.")
    st.code("""
    inventory = {'apples': 10, 'bananas': 5, 'oranges': 15}
    for fruit, count in inventory.items():
        print(f"{fruit}: {count}")
    # apples: 10
    # bananas: 5
    # oranges: 15
    
    # Only keys
    for fruit in inventory:
        print(fruit)
    """, language="python")
    
    st.subheader("Example 4: Nested Dictionaries")
    st.markdown("Complex data like configs.")
    st.code("""
    config = {
        'database': {'host': 'localhost', 'port': 5432},
        'api': {'version': '1.0', 'timeout': 30}
    }
    print(config['database']['host'])  # localhost
    """, language="python")
    
    st.subheader("Example 5: DefaultDict from collections")
    st.markdown("Auto-initialize missing keys.")
    st.code("""
    from collections import defaultdict
    
    counts = defaultdict(int)
    counts['a'] += 1
    print(counts['b'])  # 0 (default)
    print(dict(counts))  # {'a': 1, 'b': 0}
    """, language="python")
    
    st.subheader("Best Practices")
    st.markdown("""
    - Use meaningful string keys.
    - For ordered dicts, use `collections.OrderedDict` (though modern dicts preserve insertion order).
    - Avoid mutable keys (e.g., lists as keys fail).
    """)
    
    col1, col2 = st.columns(2)
    with col1:
        st.markdown("### Updating Dicts")
        st.code("""
        d1 = {'a':1}
        d2 = {'b':2}
        d1.update(d2)
        print(d1)  # {'a':1, 'b':2}
        """, language="python")
    with col2:
        st.markdown("### Time Complexity")
        st.write("All major ops: O(1) average.")

elif topic == "Sets":
    st.markdown("""
    # Sets in Python
    
    Sets are unordered, mutable collections of unique, hashable elements. Great for membership testing and mathematical set operations. Defined with `set()` or `{}` (empty must use `set()`).
    
    ## Creating Sets
    - Empty: `s = set()`
    - Elements: `colors = {'red', 'green', 'blue'}`
    - From iterable: `s = set([1, 2, 2, 3])` → `{1,2,3}`
    
    ## Set Operations
    - Add: `add(x)`, `update(iterable)` O(1) amortized
    - Remove: `remove(x)` (KeyError if missing), `discard(x)` (no error), `pop()` (arbitrary), `clear()`
    - `len(s)`, `x in s` → O(1)
    
    ## Mathematical Operations
    - Union: `s | t` or `union(t)`
    - Intersection: `s & t` or `intersection(t)`
    - Difference: `s - t` or `difference(t)`
    - Symmetric Difference: `s ^ t` or `symmetric_difference(t)`
    - Subset: `s.issubset(t)`, Superset: `s.issuperset(t)`
    
    ## Frozen Sets
    Immutable sets: `fs = frozenset(s)` – hashable, usable as dict keys.
    
    ## Examples
    """)
    
    st.subheader("Example 1: Unique Elements")
    st.markdown("Remove duplicates from a list.")
    st.code("""
    duplicates = [1, 2, 2, 3, 4, 4, 5]
    unique = set(duplicates)
    print(unique)  # {1, 2, 3, 4, 5}
    print(list(unique))  # Order not guaranteed
    """, language="python")
    
    st.subheader("Example 2: Set Operations")
    st.markdown("Venn diagram basics.")
    st.code("""
    set_a = {'apple', 'banana', 'cherry'}
    set_b = {'banana', 'orange', 'grape'}
    
    union = set_a | set_b
    print(union)  # {'apple', 'banana', 'cherry', 'orange', 'grape'}
    
    intersection = set_a & set_b
    print(intersection)  # {'banana'}
    
    difference = set_a - set_b
    print(difference)  # {'apple', 'cherry'}
    """, language="python")
    
    st.subheader("Example 3: Membership Testing")
    st.markdown("Fast lookups.")
    st.code("""
    emails = {'user1@example.com', 'user2@example.com'}
    if 'user1@example.com' in emails:
        print("User found!")  # Fast O(1)
    """, language="python")
    
    st.subheader("Example 4: Updating Sets")
    st.markdown("Adding multiple elements.")
    st.code("""
    primes = {2, 3, 5}
    primes.update([7, 11, 13])
    print(primes)  # {2, 3, 5, 7, 11, 13}
    primes.discard(4)  # No error if missing
    """, language="python")
    
    st.subheader("Example 5: Frozensets for Keys")
    st.markdown("Use in dictionaries.")
    st.code("""
    fs = frozenset([1, 2])
    data = {fs: 'pair'}
    print(data[fs])  # 'pair'
    """, language="python")
    
    st.subheader("Best Practices")
    st.markdown("""
    - Use sets for deduplication and fast lookups.
    - Avoid if order matters (use lists).
    - Elements must be hashable.
    """)
    
    st.success("Sets are perfect for 'bag of unique items' scenarios.")

    
    st.subheader("Example 1: Basic Manipulation")
    st.markdown("Cleaning and formatting text.")
    st.code("""
    text = "  hello world  "
    cleaned = text.strip().title()
    print(cleaned)  # "Hello World"
    words = cleaned.split()
    print(words)  # ['Hello', 'World']
    """, language="python")
    
    st.subheader("Example 2: F-String Formatting")
    st.markdown("Dynamic strings.")
    st.code("""
    name = "Alice"
    age = 30
    greeting = f"Hello, {name}! You are {age} years old."
    print(greeting)  # Hello, Alice! You are 30 years old.
    
    # Expressions
    pi = 3.14159
    print(f"Pi ≈ {pi:.2f}")  # Pi ≈ 3.14
    """, language="python")
    
    st.subheader("Example 3: Searching and Replacing")
    st.markdown("Text processing.")
    st.code("""
    sentence = "The quick brown fox jumps over the lazy dog."
    print(sentence.find("fox"))  # 16 (index)
    replaced = sentence.replace("fox", "cat")
    print(replaced)  # The quick brown cat jumps...
    """, language="python")
    
    st.subheader("Example 4: Splitting and Joining")
    st.markdown("CSV-like data.")
    st.code("""
    csv_data = "apple,banana,cherry"
    fruits = csv_data.split(",")
    print(fruits)  # ['apple', 'banana', 'cherry']
    
    joined = "; ".join(fruits)
    print(joined)  # apple; banana; cherry
    """, language="python")
    
    st.subheader("Example 5: Validation Methods")
    st.markdown("Input checking.")
    st.code("""
    user_input = "123abc"
    if user_input.isalnum():
        print("Alphanumeric OK")
    else:
        print("Invalid")  # Invalid
    
    if user_input[3:].isalpha():
        print("Letters after digits")
    """, language="python")
    
    st.subheader("Example 6: Multiline Strings")
    st.markdown("Docstrings or templates.")
    st.code("""
    poem = '''Roses are red,
    Violets are blue,
    Python is great,
    And so are you!'''
    print(poem)
    """, language="python")
    
    st.subheader("Best Practices")
    st.markdown("""
    - Use f-strings for modern formatting.
    - For large text, consider `str.join()` over `+`.
    - Escape special chars with `\` or raw strings.
    """)
    
    col1, col2 = st.columns(2)
    with col1:
        st.markdown("### Slicing")
        st.code("""
        s = "Python"
        print(s[::-1])  # nohtyP (reverse)
        """, language="python")
    with col2:
        st.markdown("### Escape Sequences")
        st.code("""
        print("Line1\\nLine2")  # Line1\nLine2
        """, language="python")

elif topic == "Advanced: Stacks & Queues":
    st.markdown("""
    # Advanced: Stacks & Queues
    
    Stacks and Queues are abstract data types (ADTs) implemented using lists or `collections.deque` for efficiency. Stacks follow LIFO (Last In, First Out); Queues follow FIFO (First In, First Out).
    
    ## Stacks
    - Operations: `push(x)`, `pop()`, `peek()`, `is_empty()`
    - Use cases: Function call stack, undo/redo, parsing expressions.
    
    Implementation with list (efficient for push/pop from end).
    
    ## Queues
    - Operations: `enqueue(x)`, `dequeue()`, `front()`, `is_empty()`
    - Use cases: Task scheduling, BFS in graphs, print queues.
    
    Use `deque` for O(1) ops on both ends.
    
    ## Examples
    """)
    
    st.subheader("Example 1: Stack Implementation")
    st.markdown("Simple stack class.")
    st.code("""
    class Stack:
        def __init__(self):
            self.items = []
        
        def push(self, item):
            self.items.append(item)
        
        def pop(self):
            if not self.is_empty():
                return self.items.pop()
        
        def peek(self):
            if not self.is_empty():
                return self.items[-1]
        
        def is_empty(self):
            return len(self.items) == 0
    
    # Usage
    stack = Stack()
    stack.push(1)
    stack.push(2)
    print(stack.pop())  # 2
    print(stack.peek())  # 1
    """, language="python")
    
    st.subheader("Example 2: Queue with Deque")
    st.markdown("Efficient queue.")
    st.code("""
    from collections import deque
    
    class Queue:
        def __init__(self):
            self.items = deque()
        
        def enqueue(self, item):
            self.items.append(item)
        
        def dequeue(self):
            if not self.is_empty():
                return self.items.popleft()
        
        def front(self):
            if not self.is_empty():
                return self.items[0]
        
        def is_empty(self):
            return len(self.items) == 0
    
    # Usage
    q = Queue()
    q.enqueue('task1')
    q.enqueue('task2')
    print(q.dequeue())  # task1
    print(q.front())  # task2
    """, language="python")
    
    st.subheader("Example 3: Stack for Balanced Parentheses")
    st.markdown("Validate expression.")
    st.code("""
    def is_balanced(expr):
        stack = []
        pairs = {')':'(', '}':'{', ']':'['}
        for char in expr:
            if char in pairs.values():
                stack.append(char)
            elif char in pairs:
                if not stack or stack.pop() != pairs[char]:
                    return False
        return len(stack) == 0
    
    print(is_balanced("()[]{}"))  # True
    print(is_balanced("(]"))  # False
    """, language="python")
    
    st.subheader("Example 4: Queue for BFS Simulation")
    st.markdown("Simple graph traversal simulation.")
    st.code("""
    from collections import deque
    
    def bfs(graph, start):
        visited = set()
        queue = deque([start])
        visited.add(start)
        order = []
        while queue:
            node = queue.popleft()
            order.append(node)
            for neighbor in graph[node]:
                if neighbor not in visited:
                    visited.add(neighbor)
                    queue.append(neighbor)
        return order
    
    graph = {'A': ['B', 'C'], 'B': ['D'], 'C': ['E'], 'D': [], 'E': []}
    print(bfs(graph, 'A'))  # ['A', 'B', 'C', 'D', 'E']
    """, language="python")
    
    st.subheader("Best Practices")
    st.markdown("""
    - Use `deque` for queues to avoid O(n) shifts.
    - For thread-safe stacks/queues, consider `queue` module.
    - Handle empty cases to avoid errors.
    """)
    
    st.warning("Lists are fine for stacks but inefficient for queue dequeues.")

elif topic == "Advanced: Trees & Graphs":
    st.markdown("""
    # Advanced: Trees & Graphs
    
    Trees and Graphs are hierarchical and networked data structures for modeling relationships. Trees have no cycles (acyclic graphs); Graphs can have cycles and multiple paths.
    
    ## Trees
    - Nodes with children (root at top).
    - Use cases: File systems, HTML DOM, decision trees.
    - Binary Tree: Max 2 children per node.
    - Traversal: In-order, Pre-order, Post-order, Level-order (BFS).
    
    ## Graphs
    - Nodes (vertices) connected by edges.
    - Directed/Undirected, Weighted/Unweighted.
    - Use cases: Social networks, maps, dependencies.
    - Representations: Adjacency list (dict of lists), Matrix.
    - Traversal: DFS (stack/recursion), BFS (queue).
    
    Implementations use classes or dicts.
    
    ## Examples
    """)
    
    st.subheader("Example 1: Binary Tree Node Class")
    st.markdown("Basic tree structure.")
    st.code("""
    class TreeNode:
        def __init__(self, val=0, left=None, right=None):
            self.val = val
            self.left = left
            self.right = right
    
    # Build a simple tree: 1 -> (2, 3)
    root = TreeNode(1)
    root.left = TreeNode(2)
    root.right = TreeNode(3)
    
    def inorder(node):
        if node:
            yield from inorder(node.left)
            yield node.val
            yield from inorder(node.right)
    
    print(list(inorder(root)))  # [2, 1, 3]
    """, language="python")
    
    st.subheader("Example 2: Graph with Adjacency List")
    st.markdown("Undirected graph.")
    st.code("""
    graph = {
        'A': ['B', 'C'],
        'B': ['A', 'D', 'E'],
        'C': ['A', 'F'],
        'D': ['B'],
        'E': ['B', 'F'],
        'F': ['C', 'E']
    }
    
    def print_graph(g):
        for node, neighbors in g.items():
            print(f"{node}: {neighbors}")
    
    print_graph(graph)
    # A: ['B', 'C']
    # ...
    """, language="python")
    
    st.subheader("Example 3: DFS Traversal (Recursive)")
    st.markdown("Depth-First Search.")
    st.code("""
    def dfs(graph, start, visited=None):
        if visited is None:
            visited = set()
        visited.add(start)
        print(start, end=' ')
        for neighbor in graph[start]:
            if neighbor not in visited:
                dfs(graph, neighbor, visited)
    
    dfs(graph, 'A')  # A B D E F C (one possible order)
    """, language="python")
    
    st.subheader("Example 4: Level-Order Tree Traversal (BFS)")
    st.markdown("Using queue for levels.")
    st.code("""
    from collections import deque
    
    def level_order(root):
        if not root:
            return []
        queue = deque([root])
        result = []
        while queue:
            node = queue.popleft()
            result.append(node.val)
            if node.left:
                queue.append(node.left)
            if node.right:
                queue.append(node.right)
        return result
    
    print(level_order(root))  # [1, 2, 3]
    """, language="python")
    
    st.subheader("Example 5: Detect Cycle in Graph (DFS)")
    st.markdown("Simple cycle detection.")
    st.code("""
    def has_cycle(graph, start, visited=None, rec_stack=None):
        if visited is None:
            visited = set()
        if rec_stack is None:
            rec_stack = set()
        visited.add(start)
        rec_stack.add(start)
        for neighbor in graph[start]:
            if neighbor not in visited:
                if has_cycle(graph, neighbor, visited, rec_stack):
                    return True
            elif neighbor in rec_stack:
                return True
        rec_stack.remove(start)
        return False
    
    # Add cycle: graph['F'].append('A')  # For testing
    print(has_cycle(graph, 'A'))  # False (no cycle)
    """, language="python")
    
    st.subheader("Best Practices")
    st.markdown("""
    - Use recursion for tree traversals (watch stack depth).
    - For graphs, adjacency lists for sparse graphs.
    - Libraries: `networkx` for advanced graphs (but stick to stdlib here).
    - Handle disconnected components.
    """)
    
    col1, col2 = st.columns(2)
    with col1:
        st.markdown("### Tree Height")
        st.code("""
        def height(node):
            if not node:
                return 0
            return 1 + max(height(node.left), height(node.right))
        """, language="python")
    with col2:
        st.markdown("### Graph Density")
        st.write("Edges ~ Nodes^2 / 2 for dense.")

st.markdown("</div>", unsafe_allow_html=True)